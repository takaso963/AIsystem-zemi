from io import BytesIO
from pathlib import Path
from typing import Tuple


def _decode_text(data: bytes) -> str:
    """Windowsで作った資料も読めるように、複数の文字コードを試す。"""
    for encoding in ("utf-8-sig", "utf-8", "cp932", "shift_jis"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def extract_text_from_file(file_name: str, file_bytes: bytes) -> Tuple[str, str]:
    """
    アップロードされた講義資料からテキストを取り出す。
    戻り値: (抽出テキスト, メッセージ)
    """
    suffix = Path(file_name).suffix.lower()

    if suffix in {".txt", ".md", ".csv"}:
        text = _decode_text(file_bytes)
        return text.strip(), "テキストを読み込みました。"

    if suffix == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return "", "PDFを読むための pypdf が入っていません。pip install pypdf を実行してください。"

        try:
            reader = PdfReader(BytesIO(file_bytes))
            pages = []
            for page_no, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                page_text = page_text.strip()
                if page_text:
                    pages.append(f"--- {page_no}ページ ---\n{page_text}")
            if not pages:
                return "", "PDFからテキストを抽出できませんでした。画像だけのPDFの場合はOCRが必要です。"
            return "\n\n".join(pages), f"PDF {len(reader.pages)}ページを読み込みました。"
        except Exception as e:
            return "", f"PDFの読み込みに失敗しました: {e}"

    if suffix == ".docx":
        try:
            from docx import Document
        except ImportError:
            return "", "Wordファイルを読むための python-docx が入っていません。pip install python-docx を実行してください。"

        try:
            doc = Document(BytesIO(file_bytes))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
            table_texts = []
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        table_texts.append(" | ".join(cells))
            text = "\n".join(paragraphs + table_texts)
            if not text.strip():
                return "", "Wordファイルからテキストを抽出できませんでした。"
            return text.strip(), "Wordファイルを読み込みました。"
        except Exception as e:
            return "", f"Wordファイルの読み込みに失敗しました: {e}"

    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            return "", "PowerPointを読むための python-pptx が入っていません。pip install python-pptx を実行してください。"

        try:
            prs = Presentation(BytesIO(file_bytes))
            slides = []
            for idx, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
                if texts:
                    slides.append(f"--- スライド {idx} ---\n" + "\n".join(texts))
            if not slides:
                return "", "PowerPointからテキストを抽出できませんでした。"
            return "\n\n".join(slides), f"PowerPoint {len(prs.slides)}枚を読み込みました。"
        except Exception as e:
            return "", f"PowerPointの読み込みに失敗しました: {e}"

    return "", f"未対応のファイル形式です: {suffix}"
